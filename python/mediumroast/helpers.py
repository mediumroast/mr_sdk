__version__ = '1.1'
__author__ = "Michael Hay"
__date__ = '2021-November-01'
__copyright__ = "Copyright 2022 Mediumroast, Inc. All rights reserved."


import hashlib
import datetime
import os
import json
import pathlib
import logging
import magic
import pdfx
import pydocx
import pptx
import configparser as conf
#from datetime import datetime
from geopy.geocoders import ArcGIS


class utilities:

    def __init__(self, silent=True):
        # Set up the locator for ArcGIS to get lat/long data
        self.locator = ArcGIS(timeout=2)
        self.silent = silent  # If set to True log_it will only print START, STOP and STATS

    def total_item(self, items):
        """Total items in dicts and lists and return the result.
        """
        return len(items)

    def hash_it(self, stringToHash, HASH='sha256'):
        h = hashlib.new(HASH)
        h.update(stringToHash.encode('utf-8'))
        return h.hexdigest()

    def locate(self, place):
        """Using an input string return the lat long combo using geopy
        """
        l = self.locator.geocode(place)
        return [l.longitude, l.latitude]

    def save(self, file_name, string_data):
        """ Save string content to a file
        """
        with open(file_name, 'w') as my_file:
            try:
                my_file.write(string_data)
            except IOError as err:
                return False, err
            except:
                return False, 'Something went wrong, check the file output.'
            finally:
                return True, 'Successfully wrote the data to [' + file_name + ']'

    def json_read(self, file_name):
        """Read a JSON file into memory and transform into Pythonic objects
        """
        data = ""
        with open(file_name, 'r') as my_file:
            try:
                data = my_file.read()
            except IOError as err:
                return False, err
            except:
                return False, 'Something went wrong, check the file output.'
            finally:
                return True, json.loads(data)

    def txt_read(self, file_name):
        """Read a TXT file into memory and transform into Pythonic objects
        """
        data = ""
        with open(file_name, 'r') as my_file:
            try:
                data = my_file.read()
            except IOError as err:
                return False, err
            except:
                return False, 'Something went wrong, check the file output.'
            finally:
                return True, data.strip()

    def get_item_type(self, folder, item):
        """Using libmagic return the file system object type
        """
        return magic.from_file(folder + '/' + item)

    def log_it(self, body='', function='MAIN', log_type='START'):
        """This is a very simple standin logger, ideally we will revert back to logging at some point
        """
        my_sep = ' | '
        prefix = {
            'START': '-' * 5 + '[ BEGIN: ' + function + ' ]' + '-' * 5,
            'STEP': '>' * 2 + '[' + function + ']: ',
            'STOP': '-' * 5 + '[ END: ' + function + ' ]' + '-' * 5,
            'STAT': '#' * 2 + '[' + function + ']: ',
            'DEBUG': '<< DEBUG >>>>> [' + function + ']: '
        }
        message = str()

        override = False
        if log_type == 'START':
            message = prefix[log_type]
            override = True
        elif log_type == 'STOP':
            message = prefix[log_type]
            override = True
        elif log_type == 'STAT':
            message = prefix[log_type] + body
            override = True
        elif log_type == 'STEP':
            message = prefix[log_type] + body
        elif log_type == 'DEBUG':
            message = prefix[log_type] + body

        timestamp = datetime.now()
        if override or self.silent == False:
            print(str(timestamp) + my_sep + message)

    def make_directory(self, dirname):
        """Safely create a directory and if it already exists gracefully return.
        """
        path_check = pathlib.Path(dirname)
        if not path_check.exists() and not path_check.is_dir():
            try:
                os.mkdir(dirname)
            except FileExistsError as err:
                return 0, 'The directory [' + dirname + '] already exists, unable to create.'
            except:
                return 0, 'Something abnormal happened when attempting to create the directory [' + dirname + '] please check the system logs.'
            finally:
                return 1, 'Successfully created directory [' + dirname + '].'
        else:
            logging.warning(
                'The directory [' + dirname + '] appears to already exist, no action performed.')
            return 2, 'The directory [' + dirname + '] appears to already exist, no action performed.'

    def check_file_system_object(self, full_filename):
        """A simple check to see if a file system object exists, returns false if it doesn't exist.
        """
        fso_check = pathlib.Path(full_filename)
        if fso_check.exists():
            return True, 'The file named [' + full_filename + '] already exists'
        else:
            return False, 'No such file named [' + full_filename + '] exists.'

    def correct_date(self, date_time, default_time='0000'):
        """Ensure that the date and time are correct
        """
        my_time = default_time
        my_date = date_time
        if len(date_time) > 8:
            my_time = my_date[8:]
            my_date = my_date[0:8]
        return my_date, my_time

    # TODO consider refactoring for clear separation of date and time

    # def get_date_time(self):
    #     """Get the time presently and return in two formats
    #     """
    #     the_time_is = time.localtime()
    #     time_concat = the_time_is.tm_year + the_time_is.tm_mon + \
    #         the_time_is.tm_mday + the_time_is.tm_hour + the_time_is.tm_min
    #     time_formal = time.asctime(the_time_is)
    #     return time_concat, time_formal

    def get_date_time(self):
        the_time = datetime.datetime.now()
        date = the_time.strftime("%Y%m%d") # Format YYYYMM
        time = the_time.strftime("%H%M") # Format HHMM
        return date, time


    def get_pdf_meta(self, item):
        """Obtain key metadata and raw text from PDFs. Text is suited for the purposes of named entity recognition.
        """
        pdf = pdfx.PDFx(item)
        doc_meta = pdf.get_metadata()
        return doc_meta, pdf.get_text()

    def get_docx_meta(self, item):
        """Obtain key metadata and raw text from DOCX. Text is suited for the purposes of named entity recognition.
        """
        doc_metadata = {}
        doc = pydocx.Document(item)
        properties = doc.core_properties
        doc_metadata['CreateDate'] = properties.created
        doc_metadata['type'] = properties.category
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text)
        return doc_metadata, '\n'.join(fullText)

    def get_pptx_meta(self, item):
        """Obtain key metadata and raw text from PPTX. Text is suited for the purposes of named entity recognition.
        """
        doc_metadata = {}
        preso = pptx.Presentation(item)
        properties = preso.core_properties
        doc_metadata['CreateDate'] = properties.created
        doc_metadata['type'] = properties.category
        fullText = []
        for slide in preso.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                fullText.append(shape.text)
        return doc_metadata, '\n'.join(fullText)


    def make_note(self, obj_type, creator='Mediumroast SDK load utility.'):
        """Create a sample note for an object or a child object
        """
        (time_stamp, time_string) = self.get_date_time()
        return {"1": {time_stamp: "This is an example note created for the '" + obj_type + "' object on " + time_string + " by a " + creator}}


class companies:

    def __init__(self, rewrite_config_dir="../src/mediumroast/transformers/"):
        self.RULES = {
            'dir': rewrite_config_dir,
            'company': 'company.ini',
            'study': 'study.ini',
            'interaction': 'interaction.ini'
        }
        self.rules = conf.ConfigParser()
        self.rules.read(self.RULES['dir'] + self.RULES['company'])
        self.util = utilities()

    def get_name(self, company_name):
        """Lookup a company's name from the configuration file and return it.

        As appropriate return the proper name of the company in question.  This is a helper function
        to be used as needed during the transformation process.

        Args:
            company_name (str): The company name which aligns to the name within the configuration file.

        Returns:
            string: A reformatted name of the company

        Notes:
            This initial implementation doesn't really do anything since we assume the company name is correct.
        """
        return company_name

    def get_description(self, company_name):
        """Lookup a company description from the configuration file and return it.

        As appropriate return a long form description of the company in question.  This is a helper function
        to be used as needed during the transformation process.

        Args:
            company_name (str): The company name which aligns to the name within the configuration file.

        Returns:
            string: A textual description from the configuration file OR if none is present the default.
        """
        if self.rules.has_option('descriptions', company_name):
            return self.rules.get('descriptions', company_name)
        else:
            return self.rules.get('DEFAULT', 'description')

    def get_industry(self, company_name):
        """Lookup a company industry from the configuration file and return it.

        As appropriate return the full industry of the company in question.  This is a helper function
        to be used as needed during the transformation process.

        Args:
            company_name (str): The company name which aligns to the name within the configuration file.

        Returns:
            string: A textual representation of the company's industry from the configuration file OR if none is present the default.
        """
        if self.rules.industries.get(company_name):
            return self.rules.industries[company_name]
        else:
            return self.rules.DEFAULT.industry

    def make_id(self, company_name, file_output=True):
        """Create an identifier for the company 

        Create a identifier for the company_name which is either 'NULL_GUID' or a GUID generated by hashing
        the company name with the company description.  The latter is only done when the output is to a JSON
        file.  In the implementation with the backend we should revisit this logic to see if it is enven necessary
        or perhaps the backend handles all of this.

        Args:
            company_name (str): The company name which aligns to the name within the configuration file.
            file_output (bool): A switch for determining if we're storing the output in a file or not

        Returns:
            string: A textual representation of the company's ID
        """
        description = self.get_description(company_name)
        id = 'NULL_GUID'
        if file_output:
            id = self.util.hash_it(company_name + description)
        return id


class studies:

    def __init__(self, rewrite_config_dir="../src/mediumroast/transformers/"):
        self.RULES = {
            'dir': rewrite_config_dir,
            'company': 'company.ini',
            'study': 'study.ini',
            'interaction': 'interaction.ini'
        }
        self.rules = conf.ConfigParser()
        self.rules.read(self.RULES['dir'] + self.RULES['study'])
        self.util = utilities()

    def get_name(self, study_name):
        """Lookup a study's name from the configuration file and return it.

        As appropriate return the proper name of the company in question.  This is a helper function
        to be used as needed during the transformation process.

        Args:
            study_name (str): The study name which aligns to the name within the configuration file.

        Returns:
            string: A reformatted name of the study OR the argument passed in if nothing exists in the configuration file

        """
        if self.rules.has_option('names', study_name):
            return self.rules.get('names', study_name)
        else:
            return study_name

    def get_description(self, study_name):
        """Lookup a study description from the configuration file and return it.

        As appropriate return a long form description of the study in question.  This is a helper function
        to be used as needed during the transformation process.

        Args:
            study_name (str): The study name which aligns to the name within the configuration file.

        Returns:
            string: A textual description from the configuration file OR if none is present the default.
        """
        if self.rules.has_option('descriptions', study_name):
            return self.rules.get('descriptions', study_name)
        else:
            return self.rules.get('DEFAULT', 'description')

    def make_id(self, study_name, file_output=True):
        """Create an identifier for the study 

        Create a identifier for the study_name which is either 'NULL_GUID' or a GUID generated by hashing
        the study name with the study description.  The latter is only done when the output is to a JSON
        file.  In the implementation with the backend we should revisit this logic to see if it is enven necessary
        or perhaps the backend handles all of this.

        Args:
            study_name (str): The study name which aligns to the name within the configuration file.
            file_output (bool): A switch for determining if we're storing the output in a file or not

        Returns:
            string: A textual representation of the study's ID
        """
        description = self.get_description(study_name)
        id = 'NULL_GUID'  # This should never happen, but leaving here in case something is odd in the configuration file
        if file_output:
            id = self.util.hash_it(study_name + description)
        return id


class interactions:

    def __init__(self, rewrite_config_dir="../src/mediumroast/transformers/"):
        self.RULES = {
            'dir': rewrite_config_dir,
            'company': 'company.ini',
            'study': 'study.ini',
            'interaction': 'interaction.ini'
        }
        self.rules = conf.ConfigParser()
        self.rules.read(self.RULES['dir'] + self.RULES['interaction'])
        self.util = utilities()

    def get_name(self, date, study_name, company_name):
        """Create an interaction name and return the resulting string.

        Generate an interaction name from the date and study_name

        Args:
            study_name (str): The study name which should ideally be reformatted to the proper name.
            date (str): A raw date for the interaction, this needs to be the same date fed to the interaction transform

        Returns:
            string: The generated name of the interaction which is the synthesis of the date string and study name

        """
        return str(date) + '-' + str(study_name) + '-' + str(company_name)

    def get_description(self, company_name, study_name):
        """Create a description from the interaction.

        Using a default in the configuration file merge in company and study names to generate a description for 
        the interaction.

        Args:
            company_name (str): The company name which aligns to the name within the configuration file.
            study_name (str): The study name which aligns to the name within the configuration file.

        Returns:
            string: A generated textual description generated from the company and study names.
        """
        description = self.rules.get('DEFAULT', 'description')
        description = description.replace("COMPANY", str(company_name))
        description = description.replace("STUDYNAME", str(study_name))
        return description

    def make_id(self, date, company_name, study_name, file_output=True):
        """Create an identifier for the interation.

        Create a identifier for the interaction which is either 'NULL_GUID' or a GUID generated by hashing
        the interaction name with the interaction description.  The latter is only done when the output is to a JSON
        file.  In the implementation with the backend we should revisit this logic to see if it is enven necessary
        or perhaps the backend handles all of this.

        Args:
            company_name (str): The company name which aligns to the name within the configuration file.
            study_name (str): The study name which aligns to the name within the configuration file.
            file_output (bool): A switch for determining if we're storing the output in a file or not

        Returns:
            string: A textual representation of the interactions's ID
        """
        interaction_name = self.get_name(date, study_name, company_name)
        description = self.get_description(company_name, study_name)
        id = 'NULL_GUID'  # This should never happen, but leaving here in case something is odd in the configuration file
        if file_output:
            id = self.util.hash_it(interaction_name + description)
        return id

    def get_substudy_id(self, interaction_name):
        """Lookup study and company substudy ids and return them.

        If there are rewrite rules available for the interaction name related to the substudy ids for one or both
        of the associated company and study return them else return default.  Substudy ids are needed to construct
        subcorpuses for both company and study objects to build at least proper summarizations when there is no 
        questionnaire.  While there is more research needed it is also likely a helpful grouping for KeyTheme detection
        and associated extraction.

        Args:
            interaction_name (str): The name of the interaction


        Returns:
            substudy_id (str): A textual representation of the iteration id of within the study for the interaction
            company_iteration_id (str): A textual representation of the iteration id of within the company for the interaction
        """
        substudy_id = self.rules.get('substudy_mappings', interaction_name) if self.rules.has_option(
            'substudy_mappings', interaction_name) else "default"
        company_iteration_id = self.rules.get('iterations_companies', interaction_name) if self.rules.has_option(
            'iterations_companies', interaction_name) else "default"
        return substudy_id, company_iteration_id
